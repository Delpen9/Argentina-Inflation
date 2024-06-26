from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Pt

import os
import subprocess

from scripts.cpi_analysis import (
    plot_cpi_inflation,
    get_currency_valuation_by_region,
    get_currency_valuation_by_region_excluding_argentina,
)
from scripts.household_income_analysis import plot_household_income
from scripts.poverty_levels_analysis import plot_poverty_levels
from scripts.gross_domestic_product_analysis import plot_inflation_adjusted_gdp
from scripts.fiscal_deficit_analysis import (
    plot_fiscal_deficit_as_perc_of_gdp,
    plot_cumulative_fiscal_deficit,
)
from scripts.fertility_rate_analysis import plot_fertility_rate_over_time
from scripts.exchange_rate_analysis import plot_exchange_rate_over_time
from scripts.argentina_age_analysis import plot_age_demographics_over_time
from scripts.general_population_stats_analysis import plot_general_population_statistics
from scripts.export_analysis import plot_export_percent_over_time


def generate_plots_and_collect_filenames():
    output_files = []

    output_filename = "outputs/cpi_display_lineplot.png"
    output_files.append(output_filename)
    plot_cpi_inflation(
        input_filename="inputs/cpi_inflation_data.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/currency_valuation_over_time.png"
    output_files.append(output_filename)
    get_currency_valuation_by_region(
        input_filename="inputs/cpi_inflation_data.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/currency_valuation_over_time_excluding_argentina.png"
    output_files.append(output_filename)
    get_currency_valuation_by_region_excluding_argentina(
        input_filename="inputs/cpi_inflation_data.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/household_purchasing_power_in_argentina.png"
    output_files.append(output_filename)
    plot_household_income(
        input_filename="inputs/per_capita_income_in_argentina.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/argentina_poverty_levels.png"
    output_files.append(output_filename)
    plot_poverty_levels(
        input_filename="inputs/argentina-poverty-levels.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/gross_domestic_weighted_inflation_lineplot.png"
    output_files.append(output_filename)
    plot_inflation_adjusted_gdp(
        input_filename_one="inputs/cpi_inflation_data.csv",
        input_filename_two="inputs/Gross domestic product per capita in Argentina.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/fiscal_deficit_lineplot.png"
    output_files.append(output_filename)
    plot_fiscal_deficit_as_perc_of_gdp(
        input_filename="inputs/argentina_fiscal_deficit.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/cumulative_fiscal_deficit_lineplot.png"
    output_files.append(output_filename)
    plot_cumulative_fiscal_deficit(
        input_filename="inputs/argentina_fiscal_deficit.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/fertility_rate_over_time.png"
    output_files.append(output_filename)
    plot_fertility_rate_over_time(
        input_filename="inputs/fertility_rate_data_worldwide.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/exchange_rate_over_time.png"
    output_files.append(output_filename)
    plot_exchange_rate_over_time(
        input_filename="inputs/argentine_peso_to_usd_exchange_rate.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/argentina_age_demographics_over_time.png"
    output_files.append(output_filename)
    plot_age_demographics_over_time(
        input_filename="inputs/argentina_age_demographic.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/argentina_general_population_data.png"
    output_files.append(output_filename)
    plot_general_population_statistics(
        input_filename="inputs/general_population_data.csv",
        output_file_path=output_filename,
    )

    output_filename = "outputs/export_percentage_of_gdp_over_time.png"
    output_files.append(output_filename)
    plot_export_percent_over_time(
        input_filename="inputs/historical-argentina-exports.csv",
        output_file_path=output_filename,
    )

    return output_files


def create_powerpoint(output_files: str, presentation_path: str) -> None:
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Add a title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Argentina Economic and Demographic Analysis"
    subtitle.text = "Generated Visualizations"

    # Add Bible verse slides
    bible_verses = [
        {
            "verse": "For where your treasure is, there your heart will be also.",
            "reference": "Matthew 6:21",
        },
        {
            "verse": "The rich rule over the poor, and the borrower is slave to the lender.",
            "reference": "Proverbs 22:7",
        },
        {
            "verse": "Honor the Lord with your wealth, with the firstfruits of all your crops.",
            "reference": "Proverbs 3:9",
        },
        {
            "verse": "Keep your lives free from the love of money and be content \n with what you have, because God has said, 'Never will I leave you; never will I forsake you.'",
            "reference": "Hebrews 13:5",
        },
        {
            "verse": "One person gives freely, yet gains even more; another withholds \n unduly, but comes to poverty.",
            "reference": "Proverbs 11:24",
        },
        {
            "verse": "Whoever loves money never has enough; whoever loves wealth is  \n never satisfied with their income. This too is meaningless.",
            "reference": "Ecclesiastes 5:10",
        },
    ]

    # Function to add verses to a slide
    def add_verses_slide(prs, verses):
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Remove the title placeholder if it exists
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text == "Click to add title":
                sp = shape
                slide.shapes._spTree.remove(sp._element)

        top = Inches(1)
        left = Inches(1)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(2)

        for verse in verses:
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p = text_frame.add_paragraph()
            p.text = verse["verse"]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            p.alignment = PP_ALIGN.LEFT

            p = text_frame.add_paragraph()
            p.text = verse["reference"]
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            p.alignment = PP_ALIGN.RIGHT

            top += Inches(1.5)  # Adjust the spacing between verses

    # Add the first two slides with three verses each
    add_verses_slide(prs, bible_verses[:3])
    add_verses_slide(prs, bible_verses[3:])

    # Define maximum dimensions for the images
    max_height = Inches(5.5)
    max_width = prs.slide_width - Inches(2)
    top_offset = Inches(0.25)  # Offset to lower the images slightly

    # Add slides for each PNG file
    for file in output_files:
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title

        title_text = file.split("/")[-1].replace("_", " ").replace(".png", "").title()

        # Capitalize specific terms
        title_text = title_text.replace("Gdp", "GDP").replace("Cpi", "CPI")

        title_shape.text = title_text

        # Get the image size to maintain aspect ratio
        img = Image.open(file)
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height

        # Calculate dimensions while maintaining aspect ratio
        width = max_height * aspect_ratio
        height = max_height

        if width > max_width:
            width = max_width
            height = max_width / aspect_ratio

        # Center the image horizontally and vertically
        left = (prs.slide_width - width) / 2
        top = (prs.slide_height - height) / 2 + top_offset

        pic = slide.shapes.add_picture(file, left, top, height=height, width=width)

    # Save the presentation
    prs.save(presentation_path)


def convert_pptx_to_pdf(pptx_path, pdf_path):
    # Command to convert pptx to pdf using LibreOffice
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        os.path.dirname(pdf_path),
        pptx_path,
    ]

    subprocess.run(command, check=True)

    # Rename the converted file to the specified pdf_path
    converted_pdf = os.path.join(
        os.path.dirname(pdf_path), os.path.basename(pptx_path).replace(".pptx", ".pdf")
    )
    os.rename(converted_pdf, pdf_path)

if __name__ == "__main__":
    output_files = generate_plots_and_collect_filenames()

    # Create the PowerPoint presentation
    create_powerpoint(output_files, "argentina_analysis_presentation.pptx")

    # Convert PowerPoint to PDF
    convert_pptx_to_pdf(
        "argentina_analysis_presentation.pptx", "argentina_analysis_presentation.pdf"
    )
